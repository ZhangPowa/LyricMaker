from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import textwrap
import os
import re

def process_lyrics(lyrics_text):
    # Split lines and clean lyrics
    lyrics = [line.strip() for line in lyrics_text.strip().splitlines() if line.strip()]
    lyrics = [re.sub(r'[^a-zA-Z\',]+', ' ', line).lstrip() for line in lyrics]
    return lyrics

def create_lyric_slides(songs, output_file):
    presentation = Presentation()

    # Set Tahoma font and font size
    tahoma_font = "Tahoma"
    title_font_size = Pt(30)
    lyrics_font_size = Pt(22)

    # Set the slide dimensions to widescreen (16:9)
    presentation.slide_width = 16 * 360000  # 1 inch = 360000 EMUs
    presentation.slide_height = 9 * 360000

    # Add slides for each song
    for song_title, lyrics_text in songs:
        # Process lyrics
        lyrics = process_lyrics(lyrics_text)

        # Add title slide
        title_slide_layout = presentation.slide_layouts[5]  # Index 5 corresponds to Title Slide with Content
        title_slide = presentation.slides.add_slide(title_slide_layout)

        # Set the background color of the title slide
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # RGB color for black

        # Remove the subtitle box from the title slide
        for shape in title_slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape != title_slide.shapes.title:
                shape.element.getparent().remove(shape.element)

        # Add title to the title slide
        title_box = title_slide.shapes.title
        title_box.text = song_title
        title_box.text_frame.paragraphs[0].font.name = tahoma_font
        title_box.text_frame.paragraphs[0].font.size = title_font_size
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # RGB color for white
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        left_margin = 0
        top_margin = (presentation.slide_height - title_box.height) // 2
        width = presentation.slide_width
        height = title_box.height
        title_box.left = left_margin
        title_box.top = top_margin
        title_box.width = width
        title_box.height = height
        
        # Add lyrics slides for the current song
        lines_per_slide = 4
        for i in range(0, len(lyrics), lines_per_slide):
            # Use a blank slide layout
            blank_slide_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(blank_slide_layout)

            # Set the background color of the entire slide
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # RGB color for black

            # Add a text box covering the entire width of the slide
            left_margin = 0
            top_margin = Pt(15)  # Additional space at the top (adjust as needed)
            width = presentation.slide_width
            height = presentation.slide_height - top_margin
            text_box = slide.shapes.add_textbox(left_margin, top_margin, width, height)
            text_frame = text_box.text_frame

            # Add lines to the text frame
            for j in range(min(lines_per_slide, len(lyrics) - i)):
                p = text_frame.add_paragraph()

                # Wrap the text within the specified width
                wrapped_text = textwrap.fill(lyrics[i + j], width=50)  # Adjust the width as needed
                p.text = wrapped_text

                p.font.name = tahoma_font
                p.font.size = lyrics_font_size
                p.font.color.rgb = RGBColor(255, 255, 255)  # RGB color for white

                # Center-align the text horizontally within the text box
                p.alignment = PP_ALIGN.CENTER

    # Save the presentation to the desktop
    desktop_path = os.path.join(os.path.expanduser("~"), "Desktop")
    output_file_path = os.path.join(desktop_path, output_file)
    presentation.save(output_file_path)

# Example usage
song1_lyrics = """
How I love the voice of Jesus                           
On the cross of Calvary
He declares His work is finished
He has spoken this hope to me
Though the sun had ceased its shining
Though the war appeared as lost
Christ had triumphed over evil
It was finished upon that cross


Now the curse it has been broken
Jesus paid the price for me
Full, the pardon He has offered
Great, the welcome that I receive
Boldly I approach my Father
Clothed in Jesus' righteousness
There is no more guilt to carry
It was finished upon that cross

Death was once my great opponent
Fear once had a hold on me
But the Son who died to save us
Rose that we would be free indeed!
Death was once my great opponent
Fear once had a hold on me
But the Son who died to save us
Rose that we would be free indeed!
Yes, He rose that we would be free indeed!

Free from every plan of darkness
Free to live and free to love
Death is dead and Christ is risen!
It was finished upon that cross

Onward to eternal glory
To my Saviour and my God
I rejoice in Jesus' victory
It was finished upon that cross
It was finished upon that cross
It was finished upon that cross
"""

song2_lyrics = """
Jesus, Your mercy is all my plea
I have no defense, my guilt runs too deep
The best of my works pierced Your hands and Your feet
Jesus, Your mercy is all my plea

Jesus, Your mercy is all my boast
The goodness I claim, the grounds of my hope
Whatever I lack, it's still what I need most
Jesus, Your mercy is all my boast

Praise the King who bore my sin
Took my place when I stood condemned
Oh, how good You've always been to me
I will sing of Your mercy

Jesus, Your mercy is all my rest
When fears weigh me down and enemies press
A comfort I cling to in life and in death
Jesus, Your mercy is all my rest

Praise the King who bore my sin
Took my place when I stood condemned
Oh, how good You've always been to me
I will sing of Your mercy

Jesus, Your mercy is all my joy
Forever I'll lift my heart and my voice
To sing of a treasure, no pow'r can destroy
Jesus, Your mercy is all my joy

Praise the King who bore my sin
Took my place when I stood condemned
Oh, how good You've always been to me
I will sing
Praise the King who bore my sin
Took my place when I stood condemned
Oh, how good You've always been to me
I will sing of Your mercy
"""

song3_lyrics = """
I was a wretch, I remember who I was
I was lost, I was blind
I was running out of time
Sin separated, The breach was far too wide
But from the far side of the chasm
You held me in your sight
So You made a way
Across the great divide
Left behind Heaven's throne
To build it here inside
And there at the cross, You paid the debt I owed
Broke my chains, freed my soul
For the first time I had hope

Thank you Jesus for the blood applied
Thank you Jesus, it has washed me white
Thank you Jesus, You have saved my life
Brought me from the darkness into glorious light


You took my place
Laid inside my tomb of sin
You were buried for three days
But then You walked right out again
And now death has no sting
And life has no end
For I have been transformed
By the blood of the lamb

Thank you Jesus for the blood applied
Thank you Jesus, it has washed me white
Thank you Jesus, You have saved my life
Brought me from the darkness into glorious light


There is nothing stronger
Than the wonder working power of the blood
The blood
That calls us sons and daughters
We are ransomed by our Father
Through the blood
The blood (x2)

Glory to His name
Glory to His name
There to my heart was the blood applied
Glory to His name
"""

songs_data = [
    ("Grace Alone", song1_lyrics),
    ("Come Behold the Wonderous Mystery", song2_lyrics),
    ("This is Our God", song3_lyrics)
]

create_lyric_slides(songs_data, "BOMS.pptx")
